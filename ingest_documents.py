"""
Document Ingestion Pipeline
Loads documents, creates embeddings, and stores in DocumentStore
"""

import argparse
import json
from pathlib import Path
from typing import List
from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn

from haystack import Document
from haystack.document_stores.in_memory import InMemoryDocumentStore
from haystack.components.embedders import SentenceTransformersDocumentEmbedder

from config import Config

# Document format processors
try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False

try:
    from openpyxl import load_workbook
    import xlrd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

console = Console()


class DocumentIngestionPipeline:
    """Pipeline for ingesting documents into the knowledge base"""
    
    def __init__(self):
        self.config = Config
        self.document_store = InMemoryDocumentStore()
        self.doc_embedder = None
        
    def load_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            raise ImportError("pypdf not installed. Run: pip install pypdf")
        
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    
    def load_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Run: pip install python-docx")
        
        doc = DocxDocument(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    
    def load_html(self, file_path: Path) -> str:
        """Extract text from HTML file"""
        if not HTML_AVAILABLE:
            raise ImportError("beautifulsoup4 not installed. Run: pip install beautifulsoup4")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text()
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
        return text
    
    def load_json(self, file_path: Path) -> str:
        """Extract text from JSON file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            # Convert JSON to readable text
            return json.dumps(data, indent=2)
    
    def load_excel(self, file_path: Path) -> str:
        """Extract text from Excel file (.xlsx or .xls)"""
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl/xlrd not installed. Run: pip install openpyxl xlrd")
        
        text_parts = []
        ext = file_path.suffix.lower()
        
        if ext == '.xlsx':
            # Use openpyxl for .xlsx
            wb = load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
        elif ext == '.xls':
            # Use xlrd for .xls
            wb = xlrd.open_workbook(file_path)
            for sheet in wb.sheets():
                text_parts.append(f"\n=== Sheet: {sheet.name} ===")
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_text = ' | '.join(str(cell) for cell in row if cell)
                    if row_text.strip():
                        text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    
    def load_powerpoint(self, file_path: Path) -> str:
        """Extract text from PowerPoint file"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n=== Slide {slide_num} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        
        return '\n'.join(text_parts)
    
    def load_csv(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        if not PANDAS_AVAILABLE:
            # Fallback to standard csv module
            import csv
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                return '\n'.join(' | '.join(row) for row in reader)
        
        # Use pandas for better handling
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def load_rtf(self, file_path: Path) -> str:
        """Extract text from RTF file"""
        if not RTF_AVAILABLE:
            raise ImportError("striprtf not installed. Run: pip install striprtf")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()
            return rtf_to_text(rtf_content)
    
    def load_epub(self, file_path: Path) -> str:
        """Extract text from EPUB file"""
        if not EPUB_AVAILABLE:
            raise ImportError("ebooklib not installed. Run: pip install ebooklib")
        
        book = epub.read_epub(file_path)
        text_parts = []
        
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text_parts.append(soup.get_text())
        
        return '\n\n'.join(text_parts)
    
    def load_xml(self, file_path: Path) -> str:
        """Extract text from XML file"""
        if not HTML_AVAILABLE:  # Uses lxml via BeautifulSoup
            raise ImportError("lxml not installed. Run: pip install lxml")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'lxml-xml')
            return soup.get_text()
    
    def load_documents_from_directory(self, directory: Path) -> List[Document]:
        """Load documents from a directory"""
        documents = []
        
        # Supported file extensions
        supported_extensions = {
            '.txt': 'text',
            '.md': 'text',
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.html': 'html',
            '.htm': 'html',
            '.json': 'json',
            '.xlsx': 'excel',
            '.xls': 'excel',
            '.pptx': 'powerpoint',
            '.csv': 'csv',
            '.rtf': 'rtf',
            '.epub': 'epub',
            '.xml': 'xml',
            # Code files as text
            '.py': 'text',
            '.js': 'text',
            '.java': 'text',
            '.cpp': 'text',
            '.c': 'text',
            '.h': 'text',
            '.cs': 'text',
            '.php': 'text',
            '.rb': 'text',
            '.go': 'text',
            '.rs': 'text',
            '.ts': 'text',
            '.jsx': 'text',
            '.tsx': 'text',
            '.sql': 'text',
            '.sh': 'text',
            '.yaml': 'text',
            '.yml': 'text',
            '.toml': 'text',
            '.ini': 'text',
            '.cfg': 'text',
            '.conf': 'text'
        }
        
        if not directory.exists():
            console.print(f"[yellow]WARNING: Directory not found: {directory}[/yellow]")
            console.print(f"[yellow]Creating directory: {directory}[/yellow]")
            directory.mkdir(parents=True, exist_ok=True)
            return documents
        
        # Find all supported files
        files = []
        for ext in supported_extensions.keys():
            files.extend(directory.rglob(f"*{ext}"))
        
        if not files:
            console.print(f"[yellow]WARNING: No documents found in {directory}[/yellow]")
            console.print(f"[yellow]Supported formats: {', '.join(supported_extensions.keys())}[/yellow]")
            return documents
        
        console.print(f"[cyan]Found {len(files)} document(s)[/cyan]")
        
        # Load each file
        for file_path in files:
            try:
                ext = file_path.suffix.lower()
                file_type = supported_extensions.get(ext, 'text')
                
                # Extract content based on file type
                if file_type == 'pdf':
                    content = self.load_pdf(file_path)
                elif file_type == 'docx':
                    content = self.load_docx(file_path)
                elif file_type == 'html':
                    content = self.load_html(file_path)
                elif file_type == 'json':
                    content = self.load_json(file_path)
                elif file_type == 'excel':
                    content = self.load_excel(file_path)
                elif file_type == 'powerpoint':
                    content = self.load_powerpoint(file_path)
                elif file_type == 'csv':
                    content = self.load_csv(file_path)
                elif file_type == 'rtf':
                    content = self.load_rtf(file_path)
                elif file_type == 'epub':
                    content = self.load_epub(file_path)
                elif file_type == 'xml':
                    content = self.load_xml(file_path)
                else:  # text, md, code files
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                
                if not content.strip():
                    console.print(f"[yellow]![/yellow] Skipped (empty): {file_path.name}")
                    continue
                
                # Create Haystack Document
                doc = Document(
                    content=content,
                    meta={
                        "filename": file_path.name,
                        "filepath": str(file_path),
                        "file_type": file_type,
                        "source": "local"
                    }
                )
                documents.append(doc)
                console.print(f"[green]+[/green] Loaded ({file_type}): {file_path.name}")
                
            except Exception as e:
                console.print(f"[red]-[/red] Error loading {file_path.name}: {str(e)}")
        
        return documents
    
    def load_sample_documents(self) -> List[Document]:
        """Load sample documents for testing"""
        console.print("[cyan]Loading sample documents about AI and RAG...[/cyan]")
        
        samples = [
            {
                "content": """Retrieval-Augmented Generation (RAG) is a technique that combines 
                information retrieval with text generation. It works by first retrieving relevant 
                documents from a knowledge base, then using those documents as context for a 
                language model to generate accurate, grounded responses. RAG is particularly 
                useful for question-answering systems because it allows the model to access 
                external knowledge rather than relying solely on its training data.""",
                "meta": {"filename": "rag_basics.txt", "topic": "RAG", "source": "sample"}
            },
            {
                "content": """Haystack is an open-source framework for building production-ready 
                LLM applications, retrieval-augmented generative pipelines, and state-of-the-art 
                search systems. It provides components for document stores, retrievers, readers, 
                and generators that can be composed into flexible pipelines. Haystack supports 
                various LLM providers including OpenAI, Anthropic, Cohere, and local models.""",
                "meta": {"filename": "haystack_intro.txt", "topic": "Haystack", "source": "sample"}
            },
            {
                "content": """Vector embeddings are numerical representations of text that capture 
                semantic meaning. When documents are converted to embeddings, similar documents 
                will have similar vector representations. This allows for semantic search, where 
                queries can find relevant documents based on meaning rather than just keyword 
                matching. Common embedding models include sentence-transformers and OpenAI embeddings.""",
                "meta": {"filename": "embeddings.txt", "topic": "Embeddings", "source": "sample"}
            },
            {
                "content": """Data privacy in AI systems is crucial. When using cloud-based LLMs, 
                your data is sent to external servers. For maximum privacy, consider using local 
                models like Ollama (running Llama, Mistral) or HuggingFace models that run entirely 
                on your infrastructure. Embeddings can also run locally using sentence-transformers, 
                ensuring no data leaves your system during the retrieval process.""",
                "meta": {"filename": "privacy.txt", "topic": "Privacy", "source": "sample"}
            }
        ]
        
        documents = [Document(content=doc["content"], meta=doc["meta"]) for doc in samples]
        console.print(f"[green]+[/green] Loaded {len(documents)} sample documents")
        return documents
    
    def initialize_embedder(self):
        """Initialize the document embedder"""
        console.print(f"[cyan]Initializing embedder: {self.config.EMBEDDING_MODEL}[/cyan]")
        
        with console.status("[bold cyan]Downloading embedding model..."):
            from haystack.utils import ComponentDevice
            embedder = SentenceTransformersDocumentEmbedder(
                model=self.config.EMBEDDING_MODEL,
                device=ComponentDevice.from_str("cpu")
            )
            self.doc_embedder = embedder
            self.doc_embedder.warm_up()
        
        console.print("[green]+[/green] Embedder initialized")
    
    def embed_and_store_documents(self, documents: List[Document]):
        """Create embeddings and store documents"""
        if not documents:
            console.print("[yellow]WARNING: No documents to process[/yellow]")
            return
        
        console.print(f"[cyan]Creating embeddings for {len(documents)} document(s)...[/cyan]")
        
        with console.status("[bold cyan]Generating embeddings..."):
            # Create embeddings
            docs_with_embeddings = self.doc_embedder.run(documents)
            
            # Write to document store
            self.document_store.write_documents(docs_with_embeddings["documents"])
        
        console.print(f"[green]+[/green] Stored {len(documents)} document(s) with embeddings")
    
    def save_document_store(self):
        """Save document store to disk for persistence"""
        console.print(f"[cyan]Saving document store to {self.config.DOCUMENT_STORE_PATH}...[/cyan]")
        
        # Get all documents from store
        all_docs = self.document_store.filter_documents()
        
        # Serialize to JSON
        docs_data = []
        for doc in all_docs:
            # Handle embedding - could be numpy array or list
            embedding = doc.embedding
            if embedding is not None:
                if hasattr(embedding, 'tolist'):
                    embedding = embedding.tolist()
            
            doc_dict = {
                "id": doc.id,
                "content": doc.content,
                "meta": doc.meta,
                "embedding": embedding
            }
            docs_data.append(doc_dict)
        
        # Save to file
        self.config.DOCUMENT_STORE_PATH.parent.mkdir(parents=True, exist_ok=True)
        with open(self.config.DOCUMENT_STORE_PATH, 'w', encoding='utf-8') as f:
            json.dump(docs_data, f, indent=2)
        
        console.print(f"[green]+[/green] Document store saved successfully")
    
    def run(self, source_dir: Path = None, use_samples: bool = False):
        """Run the complete ingestion pipeline"""
        console.print("\n[bold cyan]========================================[/bold cyan]")
        console.print("[bold cyan]   Document Ingestion Pipeline         [/bold cyan]")
        console.print("[bold cyan]========================================[/bold cyan]\n")
        
        # Load documents
        if use_samples:
            documents = self.load_sample_documents()
        else:
            source_dir = source_dir or self.config.DOCUMENTS_DIR
            documents = self.load_documents_from_directory(source_dir)
        
        if not documents:
            console.print("[yellow]WARNING: No documents to process. Exiting.[/yellow]")
            return
        
        # Initialize embedder
        self.initialize_embedder()
        
        # Create embeddings and store
        self.embed_and_store_documents(documents)
        
        # Save to disk
        self.save_document_store()
        
        console.print("\n[bold green]SUCCESS: Document ingestion completed![/bold green]")
        console.print(f"[green]Knowledge base ready with {len(documents)} document(s)[/green]\n")


def main():
    parser = argparse.ArgumentParser(description="Ingest documents into RAG knowledge base")
    parser.add_argument(
        "--source",
        type=str,
        help="Source directory containing documents (default: ./documents)"
    )
    parser.add_argument(
        "--samples",
        action="store_true",
        help="Load sample documents for testing"
    )
    
    args = parser.parse_args()
    
    # Display configuration
    Config.display_config()
    
    # Run pipeline
    pipeline = DocumentIngestionPipeline()
    source_dir = Path(args.source) if args.source else None
    pipeline.run(source_dir=source_dir, use_samples=args.samples)


if __name__ == "__main__":
    main()
